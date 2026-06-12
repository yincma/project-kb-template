from __future__ import annotations

from pathlib import Path
from typing import Any

from kb.parsers.ocr import OCRProcessor
from kb.parsers.registry import ParsedDocument, ParsedSection


def parse_pptx_file(path: Path, config: Any | None = None) -> ParsedDocument:
    warnings: list[str] = []
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except Exception as exc:
        return ParsedDocument(path=path, warnings=[f"Failed to parse {path}: python-pptx is unavailable: {exc}"])

    ocr_cfg = _ocr_config(config)
    office_cfg = _office_config(config)
    ocr = OCRProcessor(_cfg_value(ocr_cfg, "engine", "rapidocr")) if _cfg_value(ocr_cfg, "enabled", True) else None
    extract_images = bool(_cfg_value(office_cfg, "extract_images", True))
    extract_notes = bool(_cfg_value(office_cfg, "extract_notes", True))

    try:
        deck = Presentation(str(path))
    except Exception as exc:
        return ParsedDocument(path=path, warnings=[f"Failed to parse {path}: {exc}"])

    sections: list[ParsedSection] = []
    for slide_index, slide in enumerate(deck.slides, start=1):
        parts: list[str] = []
        title = _slide_title(slide)
        if title:
            parts.append(title)

        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                table_text = _table_text(shape.table)
                if table_text:
                    parts.append(table_text)
            elif getattr(shape, "has_text_frame", False):
                text = (shape.text or "").strip()
                if text and text != title:
                    parts.append(text)

            if extract_images and ocr and getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                try:
                    ocr_result = ocr.ocr_image(shape.image.blob)
                    if ocr_result.text.strip():
                        parts.append(f"Image OCR:\n{ocr_result.text.strip()}")
                except Exception as exc:
                    warnings.append(f"OCR failed for {path} slide {slide_index}: {exc}")

        if extract_notes:
            notes = _notes_text(slide)
            if notes:
                parts.append(f"Speaker notes:\n{notes}")

        slide_text = "\n\n".join(part for part in parts if part.strip()).strip()
        if slide_text:
            sections.append(
                ParsedSection(
                    text=slide_text,
                    heading=title or f"Slide {slide_index}",
                    slide_number=slide_index,
                    ocr_used="Image OCR:" in slide_text,
                    parser_name="pptx",
                    source_format=".pptx",
                    extraction_method="text+ocr" if "Image OCR:" in slide_text else "text",
                    asset_type="slide",
                )
            )

    return ParsedDocument(path=path, sections=sections, warnings=warnings)


def _slide_title(slide) -> str | None:
    title_shape = getattr(slide.shapes, "title", None)
    if not title_shape:
        return None
    text = (getattr(title_shape, "text", "") or "").strip()
    return text or None


def _table_text(table) -> str:
    rows: list[str] = []
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        if any(cells):
            rows.append("\t".join(cells))
    return "\n".join(rows)


def _notes_text(slide) -> str:
    try:
        notes_frame = slide.notes_slide.notes_text_frame
    except Exception:
        return ""
    text = (notes_frame.text or "").strip()
    return text


def _ocr_config(config: Any | None) -> Any | None:
    return _cfg_value(_cfg_value(config, "parsing", None), "ocr", None)


def _office_config(config: Any | None) -> Any | None:
    return _cfg_value(_cfg_value(config, "parsing", None), "office", None)


def _cfg_value(config: Any | None, name: str, default: Any = None) -> Any:
    if config is None:
        return default
    if isinstance(config, dict):
        return config.get(name, default)
    return getattr(config, name, default)
